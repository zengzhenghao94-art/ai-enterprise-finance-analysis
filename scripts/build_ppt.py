"""
Build 14-slide defense PPT — 暖金学术风格 (Warm Academic × Editorial)
Burgundy + Champagne Gold on warm cream. Solo team.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════
# COLOR SYSTEM — Warm Academic
# ═══════════════════════════════════════════
BG_CREAM      = RGBColor(0xFE, 0xFC, 0xF7)
BURGUNDY      = RGBColor(0x9B, 0x1B, 0x30)
GOLD          = RGBColor(0xC8, 0xA9, 0x51)
DARK_BROWN    = RGBColor(0x29, 0x25, 0x24)
LIGHT_GOLD    = RGBColor(0xD4, 0xC5, 0xA9)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_GOLD_BG  = RGBColor(0xF8, 0xF3, 0xE5)
DEEP_RED      = RGBColor(0x7A, 0x15, 0x25)
TEXT_GRAY     = RGBColor(0x5C, 0x55, 0x50)
CARD_BORDER   = RGBColor(0xE8, 0xDE, 0xCC)
LIGHT_CREAM   = RGBColor(0xFD, 0xFB, 0xF5)
CHART_WARM1   = RGBColor(0xC8, 0xA9, 0x51)  # gold
CHART_WARM2   = RGBColor(0x9B, 0x1B, 0x30)  # burgundy
CHART_WARM3   = RGBColor(0xD4, 0x8C, 0x5C)  # copper

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS_DIR = os.path.join(BASE_DIR, "docs")
SCREENSHOT_DASHBOARD = os.path.join(DOCS_DIR, "frontend-screenshot.png")
SCREENSHOT_NL2SQL   = os.path.join(DOCS_DIR, "nl2sql-demo.png")
SCREENSHOT_REPORT    = os.path.join(DOCS_DIR, "report-v2-preview.png")
OUTPUT_PATH = os.path.join(DOCS_DIR, "答辩PPT-暖金学术版-v2.pptx")

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Team info
TEAM_NAME = "析微"
TEAM_MEMBER = "陈昊"
SCHOOL = "江西财经大学"

# ═══════════════════════════════════════════
# FONTS — Mixed: Serif titles + Sans body
# ═══════════════════════════════════════════
FONT_TITLE = "Noto Serif SC SemiBold"
FONT_BODY  = "Microsoft YaHei Light"
FONT_MONO  = "微软雅黑"  # data tables, code labels

# ═══════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                color=DARK_BROWN, bold=False, font_name=FONT_BODY,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * 4)
    # Set East Asian font
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:altLang'), 'zh-CN')
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=Pt(16),
                          color=DARK_BROWN, font_name=FONT_BODY, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.3, bold_first=False):
    """lines: list of (text, is_bold, font_size_override, color_override) or str"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, is_bold, fs, clr = line, False, font_size, color
        else:
            text = line[0]
            is_bold = line[1] if len(line) > 1 else False
            fs = line[2] if len(line) > 2 else font_size
            clr = line[3] if len(line) > 3 else color

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = text
        p.font.size = fs
        p.font.color.rgb = clr
        p.font.bold = is_bold or (bold_first and i == 0)
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(line_spacing * 4)

    return txBox

def add_gold_line(slide, left, top, width, height=Pt(1.5), color=LIGHT_GOLD):
    """Add a thin decorative line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_burgundy_line(slide, left, top, width, height=Pt(3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BURGUNDY
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, fill_color=WHITE):
    """Add a rounded card background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = CARD_BORDER
    shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
                f"第 {num} / 14 页", font_size=Pt(9), color=TEXT_GRAY,
                alignment=PP_ALIGN.RIGHT)

def add_school_header(slide):
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(3), Inches(0.3),
                SCHOOL, font_size=Pt(9), color=TEXT_GRAY)

def add_bottom_quote(slide, text, top=Inches(6.85)):
    add_gold_line(slide, Inches(0.8), top, Inches(11.7), color=LIGHT_GOLD)
    add_textbox(slide, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.4),
                text, font_size=Pt(11), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

def add_slide_title(slide, title, subtitle=None):
    """Standard slide title with burgundy + gold underline."""
    # Gold accent line at top
    add_gold_line(slide, Inches(0), Inches(0), SLIDE_W, Pt(4), color=GOLD)
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.35), Inches(11.7), Inches(0.65),
                title, font_size=Pt(28), color=BURGUNDY, bold=True,
                font_name=FONT_TITLE)
    # Burgundy underline
    add_burgundy_line(slide, Inches(0.8), Inches(1.0), Inches(1.5))
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.4),
                    subtitle, font_size=Pt(14), color=TEXT_GRAY)

def add_table(slide, left, top, col_widths, headers, rows,
              header_bg=BURGUNDY, header_fg=WHITE, font_size=Pt(12)):
    """Add a styled table. col_widths: list of Inches, headers: list of str,
       rows: list of list of str."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(w for w in col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w,
                                          Inches(0.45) * n_rows)
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.color.rgb = header_fg
            p.font.bold = True
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER

    # Body
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            bg = WHITE if r % 2 == 0 else LIGHT_CREAM
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK_BROWN
                p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

    return table_shape

# ═══════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

# ═══ SLIDE 1: COVER ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)

# Top decorative band
add_gold_line(s, Inches(0), Inches(0), SLIDE_W, Pt(6), color=GOLD)
# Left burgundy accent bar
add_burgundy_line(s, Inches(1.2), Inches(1.8), Inches(0.06), Inches(3.8))

# Title
add_textbox(s, Inches(1.6), Inches(2.0), Inches(10.5), Inches(1.2),
            "AI 赋能企业经营分析平台", font_size=Pt(44), color=BURGUNDY,
            bold=True, font_name=FONT_TITLE)
# Subtitle
add_textbox(s, Inches(1.6), Inches(3.15), Inches(10.5), Inches(0.6),
            '从「人找数据」到「数据找人」', font_size=Pt(24), color=GOLD)

# Gold divider
add_gold_line(s, Inches(1.6), Inches(3.85), Inches(3.0), Pt(2), color=GOLD)

# Three AI icons (text-based)
icons = [
    ("🛡️  数据哨兵", "ML 自动发现异常"),
    ("🌉  数据翻译官", "自然语言追问数据"),
    ("📝  报告写手", "一键生成简报"),
]
for i, (name, desc) in enumerate(icons):
    x = Inches(1.6 + i * 3.8)
    add_textbox(s, x, Inches(4.2), Inches(3.4), Inches(0.45),
                name, font_size=Pt(16), color=BURGUNDY, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(s, x, Inches(4.6), Inches(3.4), Inches(0.35),
                desc, font_size=Pt(11), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom
add_textbox(s, Inches(1.6), Inches(6.3), Inches(10.5), Inches(0.4),
            f"{SCHOOL} · 2026 年 7 月", font_size=Pt(16), color=TEXT_GRAY,
            alignment=PP_ALIGN.LEFT)
add_page_number(s, 1)

# ═══ SLIDE 2: PAIN POINTS ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "一个月度经营分析会，准备 3 天。追问一个问题，等 20 分钟。")

# Three pain-point cards
pain_points = [
    ("🔴  发现异常", "财务月底翻报表，靠经验「感觉不对劲」", "滞后 2-4 周"),
    ("🟠  追问数据", "管理层打电话 → ERP → 导出 → 透视表", "至少 20 分钟"),
    ("🟡  写报告", "收集数据 → Excel 做表 → 手写分析 → 排 PPT", "2-3 天"),
]
for i, (title, desc, cost) in enumerate(pain_points):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.0)

    # Card bg
    add_card(s, x, y, Inches(3.8), Inches(3.2), fill_color=WHITE)

    # Title
    add_textbox(s, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.45),
                title, font_size=Pt(18), color=BURGUNDY, bold=True)
    # Desc
    add_textbox(s, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.0),
                desc, font_size=Pt(13), color=DARK_BROWN)
    # Cost highlight
    add_textbox(s, x + Inches(0.3), y + Inches(2.2), Inches(3.2), Inches(0.5),
                f"⏱ {cost}", font_size=Pt(20), color=DEEP_RED, bold=True)

# Connector arrows between cards
for i in range(2):
    x = Inches(4.7 + i * 4.1)
    add_textbox(s, x, Inches(3.3), Inches(0.3), Inches(0.4),
                "→", font_size=Pt(24), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_quote(s, "不是慢——是数据、分析、决策之间，每一段都靠人手动衔接。")
add_page_number(s, 2)

# ═══ SLIDE 3: SOLUTION ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "三个 AI 重构经营分析的三段链路")

ai_roles = [
    ("🛡️", "数据哨兵", "Isolation Forest", "ML 自动发现异常", "从「月底翻报表」\n到「实时预警」"),
    ("🌉", "数据翻译官", "通义千问", "自然语言追问数据", "从「打电话等 IT」\n到「打字 5 秒」"),
    ("📝", "报告写手", "通义千问", "一键生成简报", "从「写 3 天」\n到「3 秒初稿」"),
]
for i, (icon, name, tech, action, transform) in enumerate(ai_roles):
    x = Inches(0.6 + i * 4.2)
    y = Inches(2.0)

    add_card(s, x, y, Inches(3.9), Inches(3.8), fill_color=WHITE)

    add_textbox(s, x + Inches(0.3), y + Inches(0.2), Inches(3.3), Inches(0.5),
                f"{icon}  {name}", font_size=Pt(18), color=BURGUNDY, bold=True)

    add_textbox(s, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(0.35),
                f"技术：{tech}", font_size=Pt(11), color=GOLD)

    add_gold_line(s, x + Inches(0.3), y + Inches(1.2), Inches(3.3), Pt(1), color=LIGHT_GOLD)

    add_textbox(s, x + Inches(0.3), y + Inches(1.35), Inches(3.3), Inches(0.4),
                action, font_size=Pt(14), color=DARK_BROWN, bold=True)
    add_textbox(s, x + Inches(0.3), y + Inches(1.85), Inches(3.3), Inches(0.8),
                transform, font_size=Pt(12), color=TEXT_GRAY)

    # Bottom accent in card
    add_burgundy_line(s, x, y + Inches(3.8), Inches(3.9), Pt(3))

# Arrows between cards
for i in range(2):
    x = Inches(4.55 + i * 4.2)
    add_textbox(s, x, Inches(3.6), Inches(0.4), Inches(0.5),
                "→", font_size=Pt(28), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.5),
            "💡 NL2SQL 做了三层安全校验（意图判断 + SELECT-only + 关键字黑名单），让不懂 SQL 的管理层可以放心问——这不是技术难点，是产品意识。",
            font_size=Pt(11), color=TEXT_GRAY)
add_bottom_quote(s, "三个 AI 对应三段链路——各司其职，数据自动衔接。")
add_page_number(s, 3)

# ═══ SLIDE 4: INNOVATION 1 ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "部署到新公司，要配几天？不用配规则。", "创新点一：自适应异常检测")

# Left: Traditional approach (gray card)
add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.6), fill_color=WHITE)
add_textbox(s, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.35),
            "❌  传统规则方式", font_size=Pt(16), color=TEXT_GRAY, bold=True)
add_textbox(s, Inches(0.9), Inches(2.6), Inches(5.2), Inches(1.8),
            "• 利润率 < 10%？标红\n• 销售部毛利高、生产部毛利低\n• 同一套规则没法用\n• 换一家公司，几十条阈值全部手工重设\n• 不同部门不同数字——几十行手工配置",
            font_size=Pt(12), color=TEXT_GRAY)

# Right: Our approach (warm card)
add_card(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.6), fill_color=SOFT_GOLD_BG)
add_textbox(s, Inches(7.2), Inches(2.1), Inches(5.2), Inches(0.35),
            "✅  我们的方式：Isolation Forest", font_size=Pt(16), color=BURGUNDY, bold=True)
add_textbox(s, Inches(7.2), Inches(2.6), Inches(5.2), Inches(1.8),
            "• 不问「超过阈值了吗」\n• 问「这个数字在这个部门的历史数据里\n  看起来奇怪吗」\n• 每个部门的「正常」从数据里自己学\n• 同一套算法，各自识别各自的标准",
            font_size=Pt(12), color=DARK_BROWN)

# Bottom points
add_gold_line(s, Inches(0.8), Inches(5.0), Inches(11.7), color=LIGHT_GOLD)
add_textbox(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.5),
            "🛡️ 每个部门、每个指标，「正常」的定义从数据里自己学\n"
            "🛡️ 销售部的异常标准由销售部自己的历史数据决定——同一套算法，各自识别\n"
            "🛡️ 凌晨导入数据，早上就能开始告警。不需要为每个部门手写阈值规则——这才是真正省掉的工作量",
            font_size=Pt(13), color=DARK_BROWN)

add_bottom_quote(s, 'Isolation Forest 不问「超过阈值了吗」——它问「这个数字在这个部门的历史数据里看起来奇怪吗」。')
add_page_number(s, 4)

# ═══ SLIDE 5: INNOVATION 2 ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "报告生成 AI 知道你刚检测出了什么异常吗？", "创新点二：AI 之间的暗线——上下文自动传递机制")

# Left: Disconnected
add_card(s, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.4), fill_color=WHITE)
add_textbox(s, Inches(0.9), Inches(2.1), Inches(5.2), Inches(0.35),
            "❌  各自为政（常见做法）", font_size=Pt(16), color=TEXT_GRAY, bold=True)
add_textbox(s, Inches(0.9), Inches(2.6), Inches(5.2), Inches(1.5),
            "[异常检测 AI]   ——X——   [报告生成 AI]\n\n"
            "异常检测发现「销售部 6 月应收账款偏离 32%」\n但报告生成器不知道——它只会写\n「应收账款环比增长，需关注」",
            font_size=Pt(12), color=TEXT_GRAY)

# Right: Connected
add_card(s, Inches(6.9), Inches(2.0), Inches(5.8), Inches(2.4), fill_color=SOFT_GOLD_BG)
add_textbox(s, Inches(7.2), Inches(2.1), Inches(5.2), Inches(0.35),
            "✅  我们的方式（暗线传递）", font_size=Pt(16), color=BURGUNDY, bold=True)
add_textbox(s, Inches(7.2), Inches(2.6), Inches(5.2), Inches(1.5),
            "[异常检测] → 结构化 JSON → [报告生成 Prompt]\n\n"
            "报告写手动笔之前已经知道：\n"
            "• 本月有几个异常\n"
            "• 分别是什么级别（高/中/低风险）\n"
            "• 哪个指标哪个部门偏离了多少",
            font_size=Pt(12), color=DARK_BROWN)

# Code-like box
add_card(s, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.7), fill_color=LIGHT_CREAM)
add_textbox(s, Inches(2.8), Inches(4.9), Inches(7.7), Inches(0.5),
            "报告生成 prompt ← 部门数据 + 异常检测结果（风险等级、偏离度、影响判断）",
            font_size=Pt(12), color=BURGUNDY, bold=True, alignment=PP_ALIGN.CENTER)

add_bottom_quote(s, "不是两个 AI 各做各的——异常检测是报告生成的「眼睛」。报告写手动笔之前已经知道一切。")
add_page_number(s, 5)

# ═══ SLIDE 6: INNOVATION 3 ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "从「异常找人」到「人追问数据」到「一键出报告」", "创新点三：三段式工作流重构——完整闭环")

# Flow diagram
flow_items = [
    ("🛡️\n数据哨兵", "ML 自动\n扫描异常"),
    ("↓\n异常结果\n自动注入", ""),
    ("🌉\n数据翻译官", "自然语言\n追问异常"),
    ("↓\n追问+数据\n同时注入", ""),
    ("📝\n报告写手", "一键生成简报\n知异常·知原因"),
    ("→\n管理层", "看报告\n做决策"),
]
for i, (title, desc) in enumerate(flow_items):
    x = Inches(0.5 + i * 2.15)
    if i % 2 == 0:
        # Main node
        add_card(s, x, Inches(2.0), Inches(1.95), Inches(1.8), fill_color=WHITE)
        add_textbox(s, x + Inches(0.1), Inches(2.1), Inches(1.75), Inches(0.8),
                    title, font_size=Pt(13), color=BURGUNDY, bold=True,
                    alignment=PP_ALIGN.CENTER)
        if desc:
            add_textbox(s, x + Inches(0.1), Inches(2.85), Inches(1.75), Inches(0.6),
                        desc, font_size=Pt(9), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    else:
        # Arrow node
        add_textbox(s, x, Inches(2.5), Inches(1.8), Inches(0.8),
                    title, font_size=Pt(10), color=GOLD, bold=True,
                    alignment=PP_ALIGN.CENTER)

# Closed loop label
add_card(s, Inches(3.5), Inches(4.3), Inches(6.3), Inches(0.55), fill_color=SOFT_GOLD_BG)
add_textbox(s, Inches(3.7), Inches(4.33), Inches(5.9), Inches(0.45),
            "🔄  一个平台 · 三个 AI · 零人工衔接  ——  闭环自动流转",
            font_size=Pt(14), color=BURGUNDY, bold=True, alignment=PP_ALIGN.CENTER)

# Three-line flow
add_textbox(s, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.3),
            "[数据哨兵] ML 自动扫描异常\n"
            "       ↓ 异常结果自动注入\n"
            "[数据翻译官] 管理层用自然语言追问异常细节\n"
            "       ↓ 追问结果 + 异常数据 + 部门指标同时注入\n"
            "[报告写手] 一键生成简报——知道自己写的东西里有什么异常、为什么异常、应该怎么办",
            font_size=Pt(11), color=DARK_BROWN, alignment=PP_ALIGN.CENTER)

add_bottom_quote(s, "不是在「某个环节用了一下 AI」——是用三个 AI 重构了整条经营分析链。")
add_page_number(s, 6)

# ═══ SLIDE 7: ARCHITECTURE ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "一页讲清楚前后端 + AI 服务")

# Three-layer architecture boxes
layers = [
    ("前端  Vue 3 + Vite + Tailwind CSS + ECharts", "仪表盘 / 智能问数 / 简报展示",
     Inches(1.0), Inches(2.0), CHART_WARM1),
    ("后端  Python FastAPI + SQLAlchemy", "7 个 API 端点 + 权限校验 + 数据聚合",
     Inches(1.0), Inches(3.2), BURGUNDY),
]
# Bottom three engines (side by side)
engines = [
    ("SQLite", "3 部门 × 24 月\n6 项经营指标"),
    ("scikit-learn\nIsolation Forest", "ML 异常检测\n+ Z-score 验证"),
    ("通义千问 qwen-plus\n/ DeepSeek 备选", "NL2SQL + 数据解读\n+ 简报生成"),
]

for title, desc, left, top, accent in layers:
    add_card(s, left, top, Inches(11.3), Inches(1.0), fill_color=WHITE)
    add_burgundy_line(s, left, top, Inches(0.08), Inches(1.0))
    add_textbox(s, left + Inches(0.3), top + Inches(0.1), Inches(5.0), Inches(0.45),
                title, font_size=Pt(16), color=BURGUNDY, bold=True)
    add_textbox(s, left + Inches(5.5), top + Inches(0.1), Inches(5.5), Inches(0.45),
                desc, font_size=Pt(13), color=DARK_BROWN)

# Arrow between layers
add_textbox(s, Inches(6.2), Inches(3.05), Inches(0.5), Inches(0.3),
            "↓", font_size=Pt(18), color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Three engines at bottom
for i, (name, desc) in enumerate(engines):
    x = Inches(1.0 + i * 4.1)
    y = Inches(4.6)
    add_card(s, x, y, Inches(3.8), Inches(1.6), fill_color=WHITE if i == 0 else WHITE)
    add_textbox(s, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.6),
                name, font_size=Pt(14), color=BURGUNDY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.7),
                desc, font_size=Pt(11), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.4),
            "💡 本地 ML 模型不依赖外部 API。LLM 支持降级切换。最简情况下，NL2SQL + 简报生成仍满足「AI 应用」要求。",
            font_size=Pt(10), color=TEXT_GRAY)
add_page_number(s, 7)

# ═══ SLIDE 8: DEMO - DASHBOARD ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "经营分析驾驶舱 —— 一屏看全局", "产品演示")

# Screenshot on left
if os.path.exists(SCREENSHOT_DASHBOARD):
    pic = s.shapes.add_picture(SCREENSHOT_DASHBOARD, Inches(0.5), Inches(1.6),
                                Inches(7.5), Inches(5.0))
else:
    add_card(s, Inches(0.5), Inches(1.6), Inches(7.5), Inches(5.0), fill_color=LIGHT_CREAM)
    add_textbox(s, Inches(1.5), Inches(3.5), Inches(5.5), Inches(1.0),
                "[截图：驾驶舱仪表盘]\nfrontend-screenshot.png",
                font_size=Pt(16), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Annotation points on right
annotations = [
    ("① 6 张 KPI 卡片", "收入、成本、净利润、运营费用、\n现金流、应收账款\n实时从 SQLite 查询"),
    ("② 月度趋势折线图", "三条线（收入/成本/净利润）\n支持切换部门、年份\n数据全部联动刷新"),
    ("③ 异常预警", "🔴 高风险（红底卡片）\n🟠 中风险（橙底卡片）\n显示实际值 vs 预期范围、偏离百分比"),
]
for i, (title, desc) in enumerate(annotations):
    y = Inches(1.8 + i * 1.85)
    add_textbox(s, Inches(8.3), y, Inches(4.5), Inches(0.35),
                title, font_size=Pt(14), color=BURGUNDY, bold=True)
    add_textbox(s, Inches(8.3), y + Inches(0.35), Inches(4.5), Inches(1.2),
                desc, font_size=Pt(11), color=DARK_BROWN)
    if i < 2:
        add_gold_line(s, Inches(8.3), y + Inches(1.6), Inches(3.5), Pt(1), color=LIGHT_GOLD)

add_bottom_quote(s, "截图来自实际运行系统，不是原型图。")
add_page_number(s, 8)

# ═══ SLIDE 9: DEMO - NL2SQL + REPORT ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "左边打字问数，右边一键出报告", "产品演示")

# Left: NL2SQL screenshot
if os.path.exists(SCREENSHOT_NL2SQL):
    pic = s.shapes.add_picture(SCREENSHOT_NL2SQL, Inches(0.4), Inches(1.6),
                                Inches(5.8), Inches(3.8))
else:
    add_card(s, Inches(0.4), Inches(1.6), Inches(5.8), Inches(3.8), fill_color=LIGHT_CREAM)
    add_textbox(s, Inches(1.0), Inches(3.0), Inches(4.6), Inches(1.0),
                "[截图：智能问数]\nnl2sql-demo.png",
                font_size=Pt(14), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Right: Report screenshot
if os.path.exists(SCREENSHOT_REPORT):
    pic = s.shapes.add_picture(SCREENSHOT_REPORT, Inches(6.5), Inches(1.6),
                                Inches(6.3), Inches(3.8))
else:
    add_card(s, Inches(6.5), Inches(1.6), Inches(6.3), Inches(3.8), fill_color=LIGHT_CREAM)
    add_textbox(s, Inches(7.2), Inches(3.0), Inches(4.9), Inches(1.0),
                "[截图：经营简报]\nreport-v2-preview.png",
                font_size=Pt(14), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Labels
add_textbox(s, Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.3),
            "💬  智能问数", font_size=Pt(14), color=BURGUNDY, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(6.5), Inches(1.25), Inches(6.3), Inches(0.3),
            "📄  经营简报", font_size=Pt(14), color=BURGUNDY, bold=True,
            alignment=PP_ALIGN.CENTER)

# Features below screenshots
add_textbox(s, Inches(0.4), Inches(5.6), Inches(5.8), Inches(0.8),
            "• 自然语言 → SQL（深色代码块，可见可审查）\n"
            "• 查询结果表格 + AI 通俗解释\n"
            "• 支持追问，连续对话",
            font_size=Pt(11), color=DARK_BROWN)
add_textbox(s, Inches(6.5), Inches(5.6), Inches(6.3), Inches(0.8),
            "• 一键生成四段式简报\n"
            "  （总体概况 / 分部门分析 / 异常预警 / 行动建议）\n"
            "• 导出 ZIP：HTML 报告 + 4 张高清图表 + CSS",
            font_size=Pt(11), color=DARK_BROWN)

add_bottom_quote(s, '传统方式："打电话给财务部，等 20 分钟"。AI 方式："打字，5 秒"。')
add_page_number(s, 9)

# ═══ SLIDE 10: DATA & VALIDATION ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "验证的是框架正确性，不是数据规模")

# Four dimension cards
dims = [
    ("数据设计", "3 部门 × 6 指标 × 24 月\n含季节波动+随机噪声+增长趋势",
     "不是等差数列——\n系统能处理真实数据的 messy 性"),
    ("指标覆盖", "6 个原子指标 + 3 个派生指标\n（毛利率/净利率/毛利）",
     "NL2SQL 能动态计算派生指标\n不是简单的字段映射"),
    ("测试策略", "23 项 pytest\n覆盖 7 个 API 端点全量\n+ 前后端联调",
     "核心路径全覆盖\n边界条件有 fallback 兜底"),
    ("异常检测", "Isolation Forest\n+ Z-score 两阶段交叉验证",
     "无监督学习 + 统计检验\n双重确认，降低误报"),
]
for i, (title, what, verify) in enumerate(dims):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.0)
    add_card(s, x, y, Inches(2.95), Inches(3.0), fill_color=WHITE)
    add_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(2.55), Inches(0.35),
                title, font_size=Pt(15), color=BURGUNDY, bold=True, alignment=PP_ALIGN.CENTER)
    add_gold_line(s, x + Inches(0.5), y + Inches(0.55), Inches(1.95), Pt(1), color=GOLD)
    add_textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(2.55), Inches(1.0),
                what, font_size=Pt(10), color=DARK_BROWN, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), y + Inches(1.8), Inches(2.55), Inches(1.0),
                f"验证：{verify}", font_size=Pt(9), color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom checklist
add_card(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.0), fill_color=SOFT_GOLD_BG)
add_textbox(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.8),
            "✅ 2024/2025 两年数据切换    ✅ 降级策略完整（ML 失败→规则兜底，LLM 失败→503 优雅降级）\n"
            "✅ 前端三个模块独立渲染，一个挂了其他照常    ✅ 本地 ML 不依赖外部 API，离线可用",
            font_size=Pt(11), color=DARK_BROWN)
add_page_number(s, 10)

# ═══ SLIDE 11: VALUE PROPOSITION ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "不是在「某个环节用了一下 AI」——是三个 AI 重构了整条经营分析链")

# Main comparison table
add_table(s, Inches(0.6), Inches(2.0),
          [Inches(2.0), Inches(4.0), Inches(4.0), Inches(2.0)],
          ["环节", "传统方式", "AI 方式", "提升"],
          [
              ["异常发现", "月底翻报表，滞后 2-4 周", "ML 入库即检测，无需人工翻查", "从滞后 2-4 周到入库即知"],
              ["数据查询", "打电话→ERP→导出→透视表，20+ 分钟", "自然语言输入，5 秒出结果", "约 240x"],
              ["报告生成", "收集→做表→手写→排版，2-3 天", "AI 约 30 秒生成 + 一键导出 ZIP", "从 3 天到 30 秒"],
          ])

# Engineering depth comparison
add_textbox(s, Inches(0.6), Inches(4.0), Inches(12.3), Inches(0.4),
            "不止快——工程深度对比：", font_size=Pt(14), color=BURGUNDY, bold=True)

add_table(s, Inches(0.6), Inches(4.5),
          [Inches(2.5), Inches(4.5), Inches(5.0)],
          ["维度", "典型 AI Demo", "我们的系统"],
          [
              ["异常检测", "LLM 「帮我看看有没有问题」", "ML 无监督学习 + Z-score 交叉验证，本地运行"],
              ["数据查询安全", "LLM 直接生成 SQL 执行", "三层防线：意图判断 + SELECT-only + 关键字黑名单"],
              ["AI 故障时", "功能直接不可用", "ML 本地兜底 / LLM fallback 切换 / 前端独立渲染"],
              ["AI 协作", "各自独立，互不知情", "异常输出 → 报告输入，上下文自动传递"],
          ])

add_textbox(s, Inches(0.6), Inches(6.6), Inches(8.0), Inches(0.4),
            "让人做决策，让 AI 做剩下的。", font_size=Pt(16), color=BURGUNDY, bold=True)
add_page_number(s, 11)

# ═══ SLIDE 12: SUMMARY & OUTLOOK ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "三个 AI 重构经营分析，从「人找数据」到「数据找人」")

# Three AI review
review = [
    ("🛡️  数据哨兵", "ML 自动发现异常\n无需手工设规则\n本地模型，离线可用"),
    ("🌉  数据翻译官", "自然语言查数\n三层安全防线\n让不懂 SQL 的人放心问"),
    ("📝  报告写手", "30 秒生成简报\n异常检测结果自动注入\n知异常、知原因、知对策"),
]
for i, (name, desc) in enumerate(review):
    x = Inches(0.6 + i * 4.2)
    add_card(s, x, Inches(2.0), Inches(3.9), Inches(2.2), fill_color=WHITE)
    add_textbox(s, x + Inches(0.2), Inches(2.1), Inches(3.5), Inches(0.45),
                name, font_size=Pt(18), color=BURGUNDY, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(s, x + Inches(0.2), Inches(2.6), Inches(3.5), Inches(1.3),
                desc, font_size=Pt(12), color=DARK_BROWN, alignment=PP_ALIGN.CENTER)

# Future directions
add_gold_line(s, Inches(0.8), Inches(4.5), Inches(11.7), color=GOLD)
add_textbox(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.4),
            "未来方向", font_size=Pt(16), color=BURGUNDY, bold=True)

futures = [
    "① 接入真实 ERP 数据源（用友/金蝶 API）",
    "② 预算对比 + 预测性分析（时间序列预测下月趋势）",
    "③ SSE 流式响应——让 AI 「边想边写」",
]
for i, fut in enumerate(futures):
    add_card(s, Inches(0.6 + i * 4.2), Inches(5.1), Inches(3.9), Inches(0.55),
             fill_color=WHITE)
    add_textbox(s, Inches(0.9 + i * 4.2), Inches(5.15), Inches(3.3), Inches(0.4),
                fut, font_size=Pt(13), color=DARK_BROWN, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.4),
            "谢谢各位评委", font_size=Pt(24), color=BURGUNDY, bold=True,
            alignment=PP_ALIGN.CENTER)
add_page_number(s, 12)

# ═══ SLIDE 13: TEAM ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)

# Top gold band
add_gold_line(s, Inches(0), Inches(0), SLIDE_W, Pt(6), color=GOLD)

# Centered layout with generous whitespace
add_textbox(s, Inches(2.0), Inches(1.5), Inches(9.3), Inches(0.7),
            "关于我们", font_size=Pt(36), color=BURGUNDY, bold=True,
            alignment=PP_ALIGN.CENTER)
add_burgundy_line(s, Inches(5.7), Inches(2.3), Inches(1.9), Pt(3))

# Team name large
add_textbox(s, Inches(2.0), Inches(2.7), Inches(9.3), Inches(0.8),
            f"「{TEAM_NAME}」", font_size=Pt(32), color=GOLD, bold=True,
            alignment=PP_ALIGN.CENTER)

# Member info - solo
add_textbox(s, Inches(2.0), Inches(3.8), Inches(9.3), Inches(0.5),
            f"{TEAM_MEMBER}", font_size=Pt(22), color=DARK_BROWN,
            bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(2.0), Inches(4.4), Inches(9.3), Inches(0.4),
            "一人团队 · 全栈独立开发", font_size=Pt(14), color=TEXT_GRAY,
            alignment=PP_ALIGN.CENTER)

# School
add_gold_line(s, Inches(4.5), Inches(5.2), Inches(4.3), Pt(1), color=LIGHT_GOLD)
add_textbox(s, Inches(2.0), Inches(5.4), Inches(9.3), Inches(0.4),
            SCHOOL, font_size=Pt(18), color=BURGUNDY, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(2.0), Inches(6.3), Inches(9.3), Inches(0.4),
            "感谢各位评委的时间与指导", font_size=Pt(13), color=TEXT_GRAY,
            alignment=PP_ALIGN.CENTER)
add_page_number(s, 13)

# ═══ SLIDE 14: APPENDIX ═══
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, BG_CREAM)
add_school_header(s)
add_slide_title(s, "评委可能追问的技术细节", "附录")

qa_pairs = [
    ("Q: 数据是模拟的，真实场景能用吗？",
     "Schema 基于真实财务指标体系。切换只需改数据源连接，SQL Schema 和 Prompt 模板不变。"),
    ("Q: LLM 生成的 SQL 会不会出错？",
     "三层防护（SELECT-only + 黑名单 + try/catch）。标准财务查询准确率 > 80%。出错时返回明确错误提示。"),
    ("Q: 和 Power BI / FineBI 有什么区别？",
     "BI 工具报表是「预设」的，只能在固定范围内看。我们是「对话式」的，可以问任何问题。"
     "BI 不内置异常检测和自动报告生成。"),
    ("Q: AI API 挂了怎么办？",
     "NL2SQL 有 fallback 提示。异常检测是本地模型，不依赖外部 API。"
     "基础数据查询和可视化仪表盘完全不依赖 AI。"),
]

for i, (q, a) in enumerate(qa_pairs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.0 + row * 2.5)

    add_card(s, x, y, Inches(6.0), Inches(2.2), fill_color=WHITE)
    add_textbox(s, x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.4),
                q, font_size=Pt(14), color=BURGUNDY, bold=True)
    add_gold_line(s, x + Inches(0.2), y + Inches(0.6), Inches(5.6), Pt(1), color=LIGHT_GOLD)
    add_textbox(s, x + Inches(0.2), y + Inches(0.75), Inches(5.6), Inches(1.2),
                a, font_size=Pt(12), color=DARK_BROWN)

add_page_number(s, 14)

# ═══════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════
prs.save(OUTPUT_PATH)
print(f"[OK] PPT saved to: {OUTPUT_PATH}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Style: Warm Academic x Editorial")
print(f"   Team: {TEAM_NAME} - {TEAM_MEMBER} (solo)")
